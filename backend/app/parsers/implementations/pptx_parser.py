"""PPTX document parser using python-pptx."""

from __future__ import annotations

from pathlib import Path
from time import perf_counter
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.parsers.base import BaseParser, ParseResult, ParserInfo
from app.parsers.registry import ParserRegistry

SCHEMA_VERSION = "1.0"


class PptxParser(BaseParser):
    """Parse PPTX files into unified document blocks."""

    def info(self) -> ParserInfo:
        return ParserInfo(
            name="pptx_parser",
            supported_types=["ppt", "pptx"],
            required_gpu=False,
            required_models=[],
            version="1.0.0",
        )

    def _parse_shape(
        self,
        shape: Any,
        page: int,
        blocks: list[dict[str, Any]],
        image_index: int,
    ) -> tuple[list[dict[str, Any]], int]:
        shape_type = getattr(shape, "shape_type", None)
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                blocks, image_index = self._parse_shape(
                    child, page, blocks, image_index
                )
            return blocks, image_index

        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    blocks.append(
                        {"type": "paragraph", "text": text, "page": page}
                    )
            return blocks, image_index

        if getattr(shape, "has_table", False):
            rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
            blocks.append({"type": "table", "rows": rows, "page": page})
            return blocks, image_index

        if shape_type in (
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.LINKED_PICTURE,
        ):
            left, top = shape.left, shape.top
            width, height = shape.width, shape.height
            if all(value is not None for value in (left, top, width, height)):
                bbox = [int(left), int(top), int(left + width), int(top + height)]
            else:
                bbox = []
            image_index += 1
            blocks.append(
                {
                    "type": "image",
                    "src": f"image-{image_index}",
                    "bbox": bbox,
                    "caption": "",
                    "page": page,
                }
            )
        return blocks, image_index

    def parse(
        self,
        file_path: str,
        options: dict[str, Any] | None = None,
    ) -> ParseResult:
        started = perf_counter()
        path = Path(file_path)
        if not path.is_file():
            raise FileNotFoundError(file_path)

        presentation = Presentation(str(path))
        blocks: list[dict[str, Any]] = []
        image_index = 0

        for slide_number, slide in enumerate(presentation.slides, start=1):
            blocks.append(
                {
                    "type": "heading",
                    "level": 2,
                    "text": f"Slide {slide_number}",
                    "page": slide_number,
                }
            )
            for shape in slide.shapes:
                blocks, image_index = self._parse_shape(
                    shape, slide_number, blocks, image_index
                )

        page_count = len(presentation.slides)
        processing_time_ms = round((perf_counter() - started) * 1000, 3)
        json_data = {
            "schema_version": SCHEMA_VERSION,
            "file_type": "pptx",
            "page_count": page_count,
            "blocks": blocks,
            "meta": {
                "parser": "pptx_parser",
                "processing_time_ms": processing_time_ms,
            },
        }
        return ParseResult(
            markdown="",
            json_data=json_data,
            page_count=page_count,
            processing_time_ms=processing_time_ms,
        )


ParserRegistry.register(PptxParser())
