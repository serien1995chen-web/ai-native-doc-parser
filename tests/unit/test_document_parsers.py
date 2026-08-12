"""Unit tests for PDF, DOCX, and PPTX document parsers."""

from __future__ import annotations

import importlib
from pathlib import Path

import fitz
import pytest
from docx import Document
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.parsers.implementations import docx_parser as docx_parser_module
from app.parsers.implementations import pdf_parser as pdf_parser_module
from app.parsers.implementations import pptx_parser as pptx_parser_module
from app.parsers.implementations.docx_parser import DocxParser
from app.parsers.implementations.pdf_parser import PDFParser
from app.parsers.implementations.pptx_parser import PptxParser
from app.parsers.registry import ParserRegistry

pytestmark = pytest.mark.unit


@pytest.fixture
def pdf_path(tmp_path: Path) -> Path:
    """Create a two-page PDF with text and one embedded image."""
    path = tmp_path / "sample.pdf"
    image_path = tmp_path / "pdf_image.png"
    Image.new("RGB", (12, 12), color=(255, 0, 0)).save(image_path)

    document = fitz.open()
    for page_number in range(1, 3):
        page = document.new_page()
        page.insert_text((72, 72), f"PDF text page {page_number}")
        page.insert_text((72, 120), "Second line")
    document[0].insert_image(fitz.Rect(50, 150, 100, 200), filename=str(image_path))
    document.save(path)
    document.close()
    return path


@pytest.fixture
def docx_path(tmp_path: Path) -> Path:
    """Create a DOCX with heading, paragraphs, table, and inline image."""
    document = Document()
    document.add_heading("Doc Title", level=1)
    document.add_paragraph("Before table")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A1"
    table.cell(0, 1).text = "B1"
    table.cell(1, 0).text = "A2"
    table.cell(1, 1).text = "B2"
    document.add_paragraph("After table")

    image_path = tmp_path / "docx_image.png"
    Image.new("RGB", (8, 8), color=(0, 0, 255)).save(image_path)
    document.add_picture(str(image_path), width=Inches(1))

    path = tmp_path / "sample.docx"
    document.save(path)
    return path


@pytest.fixture
def pptx_path(tmp_path: Path) -> Path:
    """Create a two-slide PPTX with text, table, image, and grouped text."""
    presentation = Presentation()
    blank = presentation.slide_layouts[6]

    slide_one = presentation.slides.add_slide(blank)
    title_box = slide_one.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(8), Inches(1)
    )
    title_box.text = "Slide One"
    body_box = slide_one.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(1)
    )
    body_box.text = "Body one"

    table_shape = slide_one.shapes.add_table(
        2, 2, Inches(1), Inches(3), Inches(4), Inches(2)
    )
    table = table_shape.table
    table.cell(0, 0).text = "X"
    table.cell(0, 1).text = "Y"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"

    image_path = tmp_path / "pptx_image.png"
    Image.new("RGB", (8, 8), color=(0, 128, 0)).save(image_path)
    slide_one.shapes.add_picture(str(image_path), Inches(1), Inches(6), width=Inches(1))

    slide_two = presentation.slides.add_slide(blank)
    group = slide_two.shapes.add_group_shape()
    grouped_box = group.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    grouped_box.text = "Grouped text"

    path = tmp_path / "sample.pptx"
    presentation.save(path)
    return path


def _assert_common_schema(
    result,
    file_type: str,
    page_count: int,
    parser_name: str,
) -> None:
    assert result.markdown == ""
    assert result.page_count == page_count
    assert result.processing_time_ms is not None
    assert result.processing_time_ms >= 0

    json_data = result.json_data
    assert json_data["schema_version"] == "1.0"
    assert json_data["file_type"] == file_type
    assert json_data["page_count"] == page_count
    assert isinstance(json_data["blocks"], list)
    assert json_data["meta"]["parser"] == parser_name
    assert json_data["meta"]["processing_time_ms"] == result.processing_time_ms


def test_pdf_parser_returns_blocks_schema(pdf_path: Path) -> None:
    result = PDFParser().parse(str(pdf_path))
    _assert_common_schema(result, "pdf", 2, "pdf_parser")

    blocks = result.json_data["blocks"]
    block_types = [block["type"] for block in blocks]
    assert "heading" in block_types
    assert "paragraph" in block_types
    assert "image" in block_types

    headings = [block["text"] for block in blocks if block["type"] == "heading"]
    assert headings == ["Page 1", "Page 2"]
    paragraphs = [block["text"] for block in blocks if block["type"] == "paragraph"]
    assert any("PDF text page 1" in text for text in paragraphs)
    images = [block for block in blocks if block["type"] == "image"]
    assert images
    assert len(images[0]["bbox"]) == 4


def test_docx_parser_returns_blocks_schema(docx_path: Path) -> None:
    result = DocxParser().parse(str(docx_path))
    _assert_common_schema(result, "docx", 1, "docx_parser")

    blocks = result.json_data["blocks"]
    assert [block["type"] for block in blocks] == [
        "heading",
        "paragraph",
        "table",
        "paragraph",
        "image",
    ]

    heading = next(block for block in blocks if block["type"] == "heading")
    assert heading["text"] == "Doc Title"
    assert heading["level"] == 1

    table = next(block for block in blocks if block["type"] == "table")
    assert table["rows"] == [["A1", "B1"], ["A2", "B2"]]

    image = next(block for block in blocks if block["type"] == "image")
    assert image["bbox"] == []


def test_pptx_parser_returns_blocks_schema(pptx_path: Path) -> None:
    result = PptxParser().parse(str(pptx_path))
    _assert_common_schema(result, "pptx", 2, "pptx_parser")

    blocks = result.json_data["blocks"]
    headings = [block["text"] for block in blocks if block["type"] == "heading"]
    assert headings == ["Slide 1", "Slide 2"]

    paragraph_texts = [
        block["text"] for block in blocks if block["type"] == "paragraph"
    ]
    assert "Slide One" in paragraph_texts
    assert "Body one" in paragraph_texts
    assert "Grouped text" in paragraph_texts
    assert any(block["type"] == "table" for block in blocks)

    images = [block for block in blocks if block["type"] == "image"]
    assert images
    assert len(images[0]["bbox"]) == 4


def test_document_parsers_registered() -> None:
    ParserRegistry._parsers.clear()
    importlib.reload(pdf_parser_module)
    importlib.reload(docx_parser_module)
    importlib.reload(pptx_parser_module)

    assert isinstance(ParserRegistry.get_parser("pdf"), pdf_parser_module.PDFParser)
    assert isinstance(ParserRegistry.get_parser("doc"), docx_parser_module.DocxParser)
    assert isinstance(ParserRegistry.get_parser("docx"), docx_parser_module.DocxParser)
    assert isinstance(ParserRegistry.get_parser("ppt"), pptx_parser_module.PptxParser)
    assert isinstance(ParserRegistry.get_parser("pptx"), pptx_parser_module.PptxParser)


def test_parsers_raise_file_not_found(tmp_path: Path) -> None:
    missing = tmp_path / "missing-file"
    parsers = (PDFParser(), DocxParser(), PptxParser())
    for parser in parsers:
        with pytest.raises(FileNotFoundError):
            parser.parse(str(missing))


def test_empty_documents_return_valid_schema(tmp_path: Path) -> None:
    pdf_path = tmp_path / "empty.pdf"
    pdf_document = fitz.open()
    pdf_document.new_page()
    pdf_document.save(pdf_path)
    pdf_document.close()
    pdf_result = PDFParser().parse(str(pdf_path))
    _assert_common_schema(pdf_result, "pdf", 1, "pdf_parser")
    assert all(block["type"] == "heading" for block in pdf_result.json_data["blocks"])

    docx_path = tmp_path / "empty.docx"
    empty_document = Document()
    empty_document.save(docx_path)
    docx_result = DocxParser().parse(str(docx_path))
    _assert_common_schema(docx_result, "docx", 1, "docx_parser")
    assert docx_result.json_data["blocks"] == []

    pptx_path = tmp_path / "empty.pptx"
    empty_presentation = Presentation()
    empty_presentation.save(pptx_path)
    pptx_result = PptxParser().parse(str(pptx_path))
    _assert_common_schema(pptx_result, "pptx", 0, "pptx_parser")
    assert pptx_result.json_data["blocks"] == []
